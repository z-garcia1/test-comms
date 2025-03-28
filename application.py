from flask import Flask, request, render_template, jsonify, redirect, session, Response
from flask_session import Session
import boto3
import json
import os
import re
import logging
import time
import random
import string
import requests
import base64
import fitz  # PyMuPDF for PDF extraction
import pandas as pd
import pptx
from docx import Document
from werkzeug.utils import secure_filename
from bs4 import BeautifulSoup
from langchain_community.tools import DuckDuckGoSearchResults
from langchain_core.rate_limiters import InMemoryRateLimiter
import pdfplumber
from langchain.utilities.tavily_search import TavilySearchAPIWrapper
from langchain.tools.tavily_search import TavilySearchResults
from langchain_aws import ChatBedrock
import fitz  # PyMuPDF for PDF extraction
from flask_cors import CORS

#Tavily
os.environ["TAVILY_API_KEY"] = os.environ.get('SearchKey')
search = TavilySearchAPIWrapper()
tavily_tool = TavilySearchResults(api_wrapper=search)

# AWS Bedrock client setup
bedrock = boto3.client('bedrock-runtime', 
                       region_name=os.environ.get('Region'),
                       aws_access_key_id=os.environ.get('AccessKeyId'),
                       aws_secret_access_key=os.environ.get('SecretAccessKey'))
# Flask app setup
app = Flask(__name__)
app.secret_key = "your_secret_key"
CORS(app)

TI_LOGIN_URL = "https://entlogin.ti.com/as/authorization.oauth2?response_type=code&client_id=DCIT_ALL_COMMS_IR_AI&redirect_uri=https%3A%2F%2Fern2xy8fzd.us-east-1.awsapprunner.com%2Fcallback&prompt=login"
TI_USERINFO_URL = "https://entlogin.ti.com/idp/userinfo.openid"

TOKEN_CHARACTERS = string.ascii_letters + string.digits + "!?@#$&%"
VALID_TOKENS = []

def generate_secure_token():
    return "".join(random.choices(TOKEN_CHARACTERS, k=20))

def clean_expired_tokens():
    """Remove expired tokens from the list"""
    global VALID_TOKENS
    current_time = time.time()
    VALID_TOKENS = [(token, exp) for token, exp in VALID_TOKENS if exp > current_time]


# Configure upload folder
UPLOAD_FOLDER = "uploads"
ALLOWED_EXTENSIONS = {"pdf", "docx", "xlsx", "pptx", "png", "jpeg", "jpg"}
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

app.config["SESSION_TYPE"] = "filesystem"
app.config["SESSION_FILE_DIR"] = "./flask_session_files"  # optional: a directory to store session files
app.config["SESSION_PERMANENT"] = True
app.config["PERMANENT_SESSION_LIFETIME"] = 43200
Session(app)

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Disable logging

#chat_memory = []

### ✅ Page Routing ###
@app.route('/')
def home():
    user_token = request.cookies.get("token")
    clean_expired_tokens()
    if not user_token or not any(t[0] == user_token for t in VALID_TOKENS):
        return redirect("/login")
    return render_template("acknowledge.html")

PUBLIC_ROUTES = {"/chat", "/reset_chat"}

@app.route("/login")
def login():
    # Redirect the user to the TI login page
    return redirect(TI_LOGIN_URL)

@app.route("/callback")
def callback():
    """Handle the redirection back from TI login"""
    auth_code = request.args.get("code")
    if not auth_code:
        return "Authorization failed", 400
    new_token = generate_secure_token()
    expiration_time = time.time() + 28800  # Token expires in 24 hours
    VALID_TOKENS.append((new_token, expiration_time))
    token_data = token_response.json()
    access_token = token_data.get("access_token")
    user_info_response = requests.get(
        "https://entlogin.ti.com/idp/userinfo.openid",
        headers={"Authorization": f"Bearer {access_token}"}
    )
    user_info = user_info_response.json()

    print(user_info)  # You’ll find the aID or employee ID here
    session["aID"] = user_info.get("aID") or user_info.get("employee_id")
    time.sleep(1)
    response = redirect("/loading")
    response.set_cookie("token", new_token, max_age=28800)
    return response

@app.route("/loading")
def loading():
    """Loading page to store the token before redirecting to home"""
    time.sleep(2)
    return redirect("/")


@app.route("/logout")
def logout():
    """Log out and clear session"""
    session.pop("user_authenticated", None)
    response = redirect("/login")
    response.set_cookie("token", "", expires=0)
    return response

@app.route('/acknowledge', methods=["POST"])
def acknowledge():
    return render_template("index.html")

from PIL import Image
import io

def resize_image(file_path, max_size=240):
    """Resizes an image before encoding to Base64 to avoid large payloads."""
    try:
        with Image.open(file_path) as img:
            original_size = img.size  # Store original size for debugging
            img.thumbnail((max_size, max_size))  # Resize while maintaining aspect ratio

            # Save resized image to memory
            img_buffer = io.BytesIO()
            img.save(img_buffer, format=img.format)
            img_buffer.seek(0)

            # Overwrite original file with resized version
            with open(file_path, "wb") as f:
                f.write(img_buffer.getvalue())

            print(f"Resized image from {original_size} to {img.size}")
    except Exception as e:
        print(f"Error resizing image: {e}")

### ✅ File Processing Functions ###
def convert_image_to_base64(file_path):
    """Converts an image file to a Base64 string."""
    with open(file_path, "rb") as file:
        return base64.b64encode(file.read()).decode("utf-8")

def extract_text_from_pdf(file_path):
    """Extracts text from a PDF file, falling back to pdfplumber if fitz fails."""
    try:
        # Attempt extraction using fitz (PyMuPDF)
        doc = fitz.open(file_path)
        text = "\n".join(page.get_text() for page in doc).strip()
        if text:
            return text
        else:
            # Raise an exception if no text is extracted
            raise ValueError("No text extracted using fitz")
    except Exception as e:
        # Fallback to pdfplumber
        try:
            extracted_text = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        extracted_text.append(page_text)
            text = "\n".join(extracted_text).strip()
            if text:
                return text
            else:
                return f"Empty pdfplumber: {e}"
        except Exception as e2:
            return f"Both fitz and pdfplumber extraction failed: {e} | {e2}"

def extract_text_from_docx(file_path):
    """Extracts text from a Word document (.docx)."""
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text.strip()

def extract_text_from_xlsx(file_path):
    """Extracts text from an Excel file (.xlsx)."""
    dfs = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
    text = []
    for sheet_name, df in dfs.items():
        text.append(f"Sheet: {sheet_name}\n")
        text.append(df.to_string(index=False))  # Convert dataframe to string
    return "\n".join(text).strip()

def extract_text_from_pptx(file_path):
    """Extracts text from a PowerPoint file (.pptx)."""
    presentation = pptx.Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text).strip()

def process_file(file_path, file_type):
    """Processes a file and extracts its text."""

    if file_type == "pdf":
        return extract_text_from_pdf(file_path)

    elif file_type == "docx":
        return extract_text_from_docx(file_path)

    elif file_type == "xlsx":
        return extract_text_from_xlsx(file_path)

    elif file_type == "pptx":
        return extract_text_from_pptx(file_path)

    elif file_type in ["csv", "txt", "html", "odt", "rtf", "epub", "json"]:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    return "Unsupported file type."

def format_ai_response(response):
    lines = response.split("\n")
    formatted_lines = []
    for line in lines:
        match = re.match(r'^(\d+.*?:)(.*)$', line.strip())
        if match:
            bold_part = f"<b>{match.group(1)}</b>"
            remaining_part = match.group(2)
            formatted_lines.append(f"{bold_part}{remaining_part}")
        else:
            formatted_lines.append(line)
    
    return "<br>".join(formatted_lines)

import yake

def extract_keywords(text, max_keywords=10):
    # Configure YAKE for English. Adjust settings as needed.
    language = "en"
    max_ngram_size = 3
    deduplication_threshold = 0.9
    kw_extractor = yake.KeywordExtractor(lan=language,
                                           n=max_ngram_size,
                                           dedupLim=deduplication_threshold,
                                           top=max_keywords,
                                           features=None)
    keywords = kw_extractor.extract_keywords(text)
    # Returns list of tuples: (keyword, score)
    return [kw for kw, score in keywords]

def filter_history(history, dynamic_keywords):
    """
    Filters the chat history to include:
      - The last three messages.
      - Any message that contains at least one dynamically determined keyword.
    
    :param history: List of message dictionaries.
    :param dynamic_keywords: A set of keywords that you want to preserve.
    """
    selected_indexes = set()
    
    # Check each message for dynamic keywords.
    for i, msg in enumerate(history):
        message_text = msg.get("content", "")
        # Extract keywords from the current message.
        message_keywords = set(extract_keywords(message_text))
        # If any of the dynamic keywords are in the extracted keywords, select this message.
        if message_keywords.intersection(dynamic_keywords):
            selected_indexes.add(i)
    
    # Always include the last three messages regardless.
    last_one_indexes = set(range(max(0, len(history) - 1), len(history)))
    selected_indexes = sorted(selected_indexes.union(last_one_indexes))
    
    return [history[i] for i in selected_indexes]


### ✅ Chat Route (Supports Text, Files, and Web Search) ###
@app.route("/chat", methods=["POST"])
def chat():
    """Handles user messages & file uploads, allowing text-only requests as well."""
    
    chat_memory = session.get('chat_memory', [])

    # Check if the request contains JSON or form data
    if request.is_json:
        user_message = request.json.get("message", "").strip()
    else:
        user_message = request.form.get("message", "").strip()

    files = request.files.getlist("file")  # Allow multiple file uploads

    if not user_message and not files:
        return jsonify({"error": "No input provided"}), 400

    content = [{"type": "text", "text": user_message}] if user_message else []
    text_from_files = []

    for file in files:
        file_ext = file.filename.split(".")[-1].lower()

        # Ensure only one image file is uploaded at a time
        image_files = [file for file in files if file.filename.split(".")[-1].lower() in ["png", "jpeg", "jpg"]]
        if len(image_files) > 1:
            return jsonify({"error": "Only one image file can be uploaded at a time."}), 400

            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(file_path)

            try:
                # Convert image to Base64 for AI processing
                image_base64 = convert_image_to_base64(file_path)
                content.append({
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": f"image/{file_ext}",
                        "data": image_base64
                    }
                })
            finally:
                os.remove(file_path)  # Cleanup image
            continue  # Skip text processing for images

        if file_ext not in ALLOWED_EXTENSIONS:
            return jsonify({"error": f"Invalid file type: {file_ext}"}), 400

        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(file_path)

        try:
            # Extract text from supported document types
            text = process_file(file_path, file_ext)
            if text:
                text_from_files.append(text)
        finally:
            os.remove(file_path)  # Cleanup document

    # Combine extracted text from all documents
    combined_text = "\n\n".join(text_from_files)

    # Ensure text fits within Claude's context window (~200K tokens)
    if len(combined_text) > 800_000:  # Approx. 200K tokens
        combined_text = combined_text[:800_000]  # Trim excess text

    if combined_text:
        chat_memory.append({"role": "user", "content": combined_text})
    
    if combined_text:
        content.append({"type": "text", "text": combined_text})

    # Store user input in chat memory before invoking Claude
    chat_memory.append({"role": "user", "content": user_message})

    # Invoke Claude AI for processing
    ai_response = invoke_claude_bedrock(content, chat_memory)

    # Store AI response in chat memory
    chat_memory.append({"role": "assistant", "content": ai_response})

  # Format the response for display
    formatted_response = format_ai_response(ai_response)

    #quick_prompt = request.form.get("quickPrompt")
    #writing_style = data.get("writingStyle")
    print("Response: 200")
    session['chat_memory'] = chat_memory
    return jsonify({
        "response": f"""<br><br><div><pre>{formatted_response}</pre><button class="copy-button"><i class="fa-regular fa-copy"></i>&nbsp; Copy</button></div>"""
    })


@app.route("/reset_chat", methods=["POST"])
def reset_chat():
    session['chat_memory'] = []

@app.route("/chat/image", methods=["POST"])
def chat_with_image():
    """Handles image-based messages."""
    user_message = request.form.get("message", "")
    files = request.files.getlist("file")

    if not files or len(files) != 1:
        return jsonify({"error": "Only one image can be uploaded at a time."}), 400

    file = files[0]
    file_ext = file.filename.split(".")[-1].lower()

    if file_ext not in ["png", "jpeg", "jpg"]:
        return jsonify({"error": "Invalid image format. Use PNG, JPEG, or JPG."}), 400

    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    file.save(file_path)

    try:
        ai_response = invoke_claude_with_image(file_path, file_ext, user_message)
    finally:
        os.remove(file_path)  # Clean up

    formatted_response = format_ai_response(ai_response)
    return jsonify({"response": f"<br><br><div><pre>{formatted_response}</pre><button class='copy-button'><i class='fa-regular fa-copy'></i>&nbsp; Copy</button></div>"})

def get_text_from_content(content):
    if isinstance(content, list):
        return " ".join(item.get("text", "") for item in content if isinstance(item, dict))
    return content

### ✅ Claude AI Invocation ###
def invoke_claude_bedrock(content, chat_memory):
    """Sends text-based content to Claude AI via AWS Bedrock, preserving chat history."""

    # Append the new user message (without system prompts) to form full history.
    full_history = chat_memory + [{"role": "user", "content": content}]
    
    # Dynamically determine the important keywords.
    # For instance, extract keywords from the entire conversation:
    # Concatenate all text from the full history using our helper function.
    all_text = " ".join(get_text_from_content(msg.get("content", "")) for msg in full_history)
    dynamic_keywords = set(extract_keywords(all_text))

    filtered_history = filter_history(full_history, dynamic_keywords)

    messages = filtered_history

    payload = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 4000,
        "messages": messages  # Include full chat history
    }

    response = bedrock.invoke_model(
        modelId="anthropic.claude-3-5-sonnet-20240620-v1:0",
        contentType="application/json",
        accept="application/json",
        body=json.dumps(payload)
    )

    response_body = response["body"].read().decode("utf-8")
    result = json.loads(response_body)

    if "content" in result and isinstance(result["content"], list):
        extracted_text = "\n".join(item["text"] for item in result["content"] if item["type"] == "text")
    else:
        extracted_text = "No valid response from Claude."

    return extracted_text


def invoke_claude_with_image(file_path, file_ext, user_message):
    """Handles image-based requests to Claude 3.5 Sonnet."""

    resize_image(file_path, max_size=240)
  
    with open(file_path, "rb") as file:
        base64_string = base64.b64encode(file.read()).decode("utf-8")

    payload = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 1024,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": f"image/{file_ext}",
                            "data": base64_string
                        }
                    },
                    {
                        "type": "text",
                        "text": user_message if user_message else "Analyze this image."
                    }
                ]
            }
        ]
    }

    response = bedrock.invoke_model(
        modelId="anthropic.claude-3-5-sonnet-20240620-v1:0",
        contentType="application/json",
        accept="application/json",
        body=json.dumps(payload)
    )

    response_body = response["body"].read().decode("utf-8")
    result = json.loads(response_body)

    if "content" in result and isinstance(result["content"], list):
        extracted_text = "\n".join(item["text"] for item in result["content"] if item["type"] == "text")
    else:
        extracted_text = "No valid response from Claude."

    return extracted_text


### ✅ Web Search Setup ###
rate_limiter = InMemoryRateLimiter(requests_per_second=0.2, check_every_n_seconds=0.1)
search_tool = DuckDuckGoSearchResults(rate_limiter=rate_limiter)

USER_AGENTS = [
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/112.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/110.0.0.0 Safari/537.36",
    "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/108.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Edge/112.0.0.0"
]

def safe_search(query):
    """Search DuckDuckGo while rotating User-Agent."""
    time.sleep(random.uniform(3, 10))
    response = search_tool.invoke(query)
    return response

@app.route("/export-chat-txt", methods=["GET"])
def export_chat_txt():
    """Generates a TXT file of chat history and returns it as a downloadable file."""

    chat_memory = session.get("chat_memory", [])

    if not chat_memory:
        return "No chat history found", 404

    chat_json = json.dumps(chat_memory)

    # Send the file as a downloadable response
    return Response(
        chat_json,
        mimetype="text/plain",
        headers={"Content-Disposition": "attachment; filename=chat_history.txt"}
    )

@app.route("/upload-chat-txt", methods=["POST"])
def upload_chat_txt():
    """Allows users to upload a chat_history.txt file and restore chat_memory to the chat box."""

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]

    # Ensure it's a text file
    if not file.filename.endswith(".txt"):
        return jsonify({"error": "Invalid file type. Please upload a .txt file"}), 400

    try:
        # Read the file content
        file_content = file.read().decode("utf-8")
        
        # Convert JSON text back to Python list
        restored_chat = json.loads(file_content)

        # Save restored chat in session
        session["chat_memory"] = restored_chat

        # Return chat history as JSON for the frontend
        return jsonify({"success": "Chat history restored successfully!", "chat_memory": restored_chat}), 200

    except Exception as e:
        return jsonify({"error": f"Failed to process file: {str(e)}"}), 500

LAMBDA_API_URL = "https://g4glmetryxjabzxbkv5sl5ksru0yawtx.lambda-url.us-east-1.on.aws/"

def call_lambda(pdf_data):
    response = requests.post(LAMBDA_API_URL, json={"body": base64.b64encode(pdf_data).decode()}, timeout=30)
    return response.json()

def call_bedrock(image_b64):
    payload = {
        "modelId": "meta.llama3-2-11b-instruct-v1:0",
        "inputText": "",
        "inputImage": image_b64,
        "parameters": {"maxLength": 512}
    }
    response = bedrock.invoke_model(
        modelId = "meta.llama3-2-11b-instruct-v1:0",
        contentType = "application/json",
        accept = "application/json",
        body = json.dumps(payload)
    )
    return json.loads(response["body"].read())["text"]

@app.route("/upload-pdf-to-image", methods=["POST"])
def upload_pdf():
    if "file" not in request.files:
        return jsonify({"error": "No file"}), 400
    
    file = request.files["file"]
    pdf_data = file.read()

    lambda_response = call_lambda(pdf_data)
    if "error" in lambda_response:
        return jsonify({"error": "Error"}), 500
    
    markdown_output = []
    for image_b64 in lambda_response["images"]:
        extracted_text = call_bedrock(image_b64)
        markdown_output.append(f"- {extracted_text}")

    return jsonify({"markdown": "\n".join(markdown_output)})

LAMBDA_URL_WEB = "https://iauc34s2dgg66w4oxlb7wfr5d40dxgqp.lambda-url.us-east-1.on.aws/"

@app.route("/search-agent", methods=["POST"])
def search_agent():
    data = request.get_json(silent=True)
    if not data or "query" not in data:
        return jsonify({"error": "Query cannot be empty"}), 400
    query = data["query"].strip()
    if not query:
        return jsonify({"error": "Empty query"})
    chat_memory = session.get('chat_memory', [])
    response = requests.post(LAMBDA_URL_WEB, json={"query": query})
    lambda_response = response.json()
    ai_response = lambda_response.get("ai_response", "No response provided.")
    chat_memory.append({"role": "assistant", "content": ai_response})
    session['chat_memory'] = chat_memory
    return jsonify(lambda_response)

### ✅ Flask App Execution for AWS App Runner ###
if __name__ == "__main__":
    app.run(debug=True)
